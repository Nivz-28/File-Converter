import os
from docx import Document
from pptx import Presentation
from io import StringIO
from pdfminer.converter import TextConverter
from pdfminer.layout import LAParams
from pdfminer.pdfdocument import PDFDocument
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.pdfpage import PDFPage
from pdfminer.pdfparser import PDFParser

def clean_text(text):
    # Replace non-breaking spaces with regular spaces, remove form feeds, soft hyphens, and the string [NBSP]
    cleaned = text.replace('\xa0', ' ').replace('\x0c', '').replace('\xad', '').replace('[NBSP]', ' ')
    return cleaned

def convert_file(file_path, output_path):
    file_name = os.path.basename(file_path)
    file_ext = os.path.splitext(file_name)[1]

    if file_ext == ".pdf":
        output_string = StringIO()
        with open(file_path, 'rb') as f:
            parser = PDFParser(f)
            doc = PDFDocument(parser)
            rsrcmgr = PDFResourceManager()
            device = TextConverter(rsrcmgr, output_string, laparams=LAParams())
            interpreter = PDFPageInterpreter(rsrcmgr, device)
            for page in PDFPage.create_pages(doc):
                interpreter.process_page(page)
        text = output_string.getvalue()

    elif file_ext == ".docx":
        doc = Document(file_path)
        text = "\n".join([para.text for para in doc.paragraphs])

    elif file_ext == ".pptx":
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"

    else:
        raise ValueError(f"Unsupported file type: {file_ext}")

    text = clean_text(text)

    with open(output_path, 'w') as f:
        f.write(text)

def main():
   file_path = "/content/random.pdf"
   output_path = "/content/outputpdf1.txt"
   convert_file(file_path, output_path)

if __name__ == "__main__":
   main()